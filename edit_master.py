# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

FONT = "Elms Sans"
DARKGREEN = "012D2B"
TEAL = "0C6D61"
MINT = "A3CCAB"
ORANGE = "F26800"
LIGHT = "EAF0EC"
GRAY = "5F6369"
WHITE = "FFFFFF"
CARD_DARK = "0A3A37"

EMU_IN = 914400


def IN(v):
    return Emu(int(round(v * EMU_IN)))


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


def set_spc(run, val):
    run.font._rPr.set("spc", str(val))


def set_font(run, size=18, bold=False, color=DARKGREEN, italic=False, spc=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    # ensure ea/cs typefaces too (mirrors master convention)
    rPr = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)
    if spc is not None:
        set_spc(run, spc)


FOOTER_ZONE_IN = 8.8  # shapes below this y are footer (logo, divider, footer texts) — always keep


def clear_shapes(slide, keep_names):
    """Remove all shapes except the corner logo and anything sitting in the footer zone
    (divider line, 'Ekotruck — Deck mestre' / 'TEMPLATE — ...' texts)."""
    to_remove = []
    for shp in slide.shapes:
        if shp.name in keep_names:
            continue
        if shp.top is not None and shp.top >= IN(FOOTER_ZONE_IN):
            continue
        to_remove.append(shp)
    for shp in to_remove:
        shp._element.getparent().remove(shp._element)


def add_rect(slide, x, y, w, h, fill, rounded=False, radius=0.03):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, IN(x), IN(y), IN(w), IN(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    if rounded:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_oval(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(d), IN(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None, space_after=None):
    """runs: list of paragraphs, each a list of (text, size, bold, color, italic) tuples,
    or a flat list of tuples for a single paragraph."""
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(25400)
    tf.margin_right = Emu(25400)
    tf.margin_top = Emu(25400)
    tf.margin_bottom = Emu(25400)
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for text, size, bold, color, italic in para_runs:
            r = p.add_run()
            r.text = text
            set_font(r, size=size, bold=bold, color=color, italic=italic)
    return tb


def add_bullet_list(slide, x, y, w, h, items, size=14, color=DARKGREEN, space_after=8):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(25400)
    tf.margin_right = Emu(25400)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        pPr = p._pPr
        if pPr is None:
            pPr = p.get_or_add_pPr()
        pPr.set("marL", "228600")
        pPr.set("indent", "-228600")
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pPr.append(buFont)
        pPr.append(buChar)
        r = p.add_run()
        r.text = item
        set_font(r, size=size, bold=False, color=color)
    return tb


def add_icon(slide, x, y, w, h, path):
    return slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))


def icon_circle(slide, x, y, d, circle_color, icon_path):
    add_oval(slide, x, y, d, circle_color)
    pad = d * 0.28
    add_icon(slide, x + pad / 2, y + pad / 2, d - pad, d - pad, icon_path)


def add_arrow_connector(slide, x1, y1, x2, y2, color=MINT, width=2.25):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def set_footer_label(slide, text):
    """Find the footer 'TEMPLATE — ...' textbox (last shape with an orange or teal
    run starting with 'TEMPLATE') and replace its text."""
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        full = "".join(r.text for p in shp.text_frame.paragraphs for r in p.runs)
        if full.startswith("TEMPLATE"):
            p = shp.text_frame.paragraphs[0]
            for r in p.runs:
                r.text = ""
            p.runs[0].text = text
            return True
    return False


ICON_DIR = r"C:\Users\andre.eiras\Desktop\Apresentações Ek\build\assets"


def icon(name, color):
    return f"{ICON_DIR}\\{name}_{color}.png"


prs = Presentation(r"C:\Users\andre.eiras\Desktop\Apresentações Ek\build\master_stage1.pptx")

slides = prs.slides
s36 = slides[30]  # Ponto de Decisão — duas opções (dup of slide4)
s37 = slides[31]  # Ponto de Decisão — lista consolidada (dup of slide4)
s38 = slides[32]  # Fluxograma com bifurcação (dup of slide21)

# ============================================================
# Slide 36 — TEMPLATE: Ponto de decisão (duas opções)
# ============================================================
keep = {"Image 0"}
clear_shapes(s36, keep)

# badge
badge = add_rect(s36, 1.147, 0.955, 3.6, 0.58, ORANGE, rounded=True, radius=0.5)
add_text(s36, 1.147, 0.955, 3.6, 0.58,
         [("PONTO DE DECISÃO", 14.5, True, WHITE, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s36, 1.147, 1.75, 15.5, 1.05,
         [("Pergunta que precisa de decisão da diretoria", 40, True, WHITE, False)])

add_text(s36, 1.147, 2.95, 17.5, 1.1,
         [("Contexto objetivo do que está em jogo — por que essa decisão é necessária agora e "
           "o que muda dependendo do caminho escolhido.", 16, False, MINT, False)],
         line_spacing=1.25)

opts = [
    ("Opção A", "Descrição objetiva do primeiro caminho possível e o que ele exige.", "Principal vantagem desta opção"),
    ("Opção B", "Descrição objetiva do segundo caminho possível e o que ele exige.", "Principal vantagem desta opção"),
]
cw, cx0, gap, cy, ch = 8.85, 1.147, 0.5, 4.35, 3.7
for i, (t, d, pro) in enumerate(opts):
    x = cx0 + i * (cw + gap)
    card = add_rect(s36, x, cy, cw, ch, CARD_DARK, rounded=True, radius=0.03)
    accent = add_rect(s36, x, cy, cw, 0.07, ORANGE)
    add_text(s36, x + 0.5, cy + 0.5, cw - 1.0, 0.8, [(t, 24, True, WHITE, False)])
    add_text(s36, x + 0.5, cy + 1.4, cw - 1.0, 1.1, [(d, 16, False, LIGHT, False)], line_spacing=1.25)
    add_text(s36, x + 0.5, cy + 2.75, cw - 1.0, 0.7, [("\u21b3 " + pro, 16.5, True, ORANGE, False)])

add_text(s36, 1.147, 8.9, 17.5, 0.55,
         [("Recomendação: levar como ponto de discussão à diretoria — não como definição já fechada.",
           15, False, MINT, True)])

set_footer_label(s36, "TEMPLATE — Ponto de decisão (duas opções)")

notes = s36.notes_slide
notes.notes_text_frame.text = ("TEMPLATE - Ponto de decisão (duas opções). Fundo escuro com badge laranja "
                                "'PONTO DE DECISÃO', contexto e dois caminhos possíveis lado a lado, cada "
                                "um com a vantagem principal em destaque. Use quando a reunião precisa "
                                "parar em um ponto específico para a diretoria escolher entre dois caminhos.")

# ============================================================
# Slide 37 — TEMPLATE: Ponto de decisão (lista consolidada)
# ============================================================
clear_shapes(s37, keep)

badge2 = add_rect(s37, 1.147, 0.955, 3.6, 0.58, ORANGE, rounded=True, radius=0.5)
add_text(s37, 1.147, 0.955, 3.6, 0.58,
         [("PONTO DE DECISÃO", 14.5, True, WHITE, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s37, 1.147, 1.75, 15.5, 1.05,
         [("Pontos de decisão para a diretoria", 40, True, WHITE, False)])

points = [
    ("bi_desmonte", "Área 1", "Do que se trata essa decisão, em uma frase objetiva."),
    ("ci_document", "Área 2", "Do que se trata essa decisão, em uma frase objetiva."),
    ("bi_armazenamento", "Área 3", "Do que se trata essa decisão, em uma frase objetiva."),
    ("bi_motores", "Área 4", "Do que se trata essa decisão, em uma frase objetiva."),
]
rw, rx, ry0, rh, rgap = 17.7, 1.147, 2.85, 1.5, 0.2
for i, (ic, t, d) in enumerate(points):
    y = ry0 + i * (rh + rgap)
    add_rect(s37, rx, y, rw, rh, CARD_DARK, rounded=True, radius=0.06)
    icon_circle(s37, rx + 0.4, y + rh / 2 - 0.44, 0.88, ORANGE, icon(ic, "white"))
    add_text(s37, rx + 1.6, y + 0.16, 3.4, rh - 0.32, [(t, 19, True, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s37, rx + 5.1, y + 0.16, rw - 5.5, rh - 0.32, [(d, 15.5, False, LIGHT, False)],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

set_footer_label(s37, "TEMPLATE — Ponto de decisão (lista consolidada)")

notes2 = s37.notes_slide
notes2.notes_text_frame.text = ("TEMPLATE - Ponto de decisão (lista consolidada). Fundo escuro, badge "
                                 "laranja e uma linha por decisão pendente (ícone + área + do que se "
                                 "trata). Use para recapitular, no fechamento de uma reunião, todos os "
                                 "pontos que ficaram pendentes de aprovação da diretoria.")

# ============================================================
# Slide 38 — TEMPLATE: Fluxograma com bifurcação e convergência
# ============================================================
keep38 = {"Image 0"}
clear_shapes(s38, keep38)

add_text(s38, 1.147, 0.955, 6, 0.42, [("PROCESSO", 19.5, True, TEAL, False)], )
add_text(s38, 1.147, 1.5, 15.5, 1.05,
         [("Como o processo se ramifica e converge", 40, True, DARKGREEN, False)])

# step 1
add_rect(s38, 6.57, 2.85, 8.75, 1.05, DARKGREEN, rounded=True, radius=0.1)
add_text(s38, 7.0, 2.98, 8.0, 0.42, [("1. Etapa inicial", 16.5, True, WHITE, False)])
add_text(s38, 7.0, 3.4, 8.0, 0.45, [("Descrição breve do primeiro passo do processo.", 12.5, False, MINT, False)])

add_arrow_connector(s38, 10.95, 3.9, 10.95, 4.25, color=TEAL)

flows = [
    ("Fluxo A", "O que acontece nesta etapa em paralelo."),
    ("Fluxo B", "O que acontece nesta etapa em paralelo."),
    ("Fluxo C", "O que acontece nesta etapa em paralelo."),
    ("Fluxo D", "O que acontece nesta etapa em paralelo."),
]
fw, fx0, fgap, fy, fh = 4.35, 1.147, 0.35, 4.55, 2.15
centers = [fx0 + fw / 2 + i * (fw + fgap) for i in range(4)]
# branch line
branch = slides = None
conn = s38.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(centers[0]), IN(4.25), IN(centers[3]), IN(4.25))
conn.line.color.rgb = rgb(MINT)
conn.line.width = Pt(2.25)
for cx in centers:
    add_arrow_connector(s38, cx, 4.25, cx, fy, color=MINT)

for i, (t, d) in enumerate(flows):
    x = fx0 + i * (fw + fgap)
    add_rect(s38, x, fy, fw, fh, LIGHT, rounded=True, radius=0.05)
    add_text(s38, x + 0.28, fy + 0.28, fw - 0.56, 0.42, [(t, 15.5, True, DARKGREEN, False)])
    add_text(s38, x + 0.28, fy + 0.75, fw - 0.56, fh - 0.95, [(d, 12, False, GRAY, False)], line_spacing=1.15)

converge_y = fy + fh + 0.35
for cx in centers:
    add_arrow_connector(s38, cx, fy + fh, cx, converge_y, color=MINT)
conn2 = s38.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(centers[0]), IN(converge_y), IN(centers[3]), IN(converge_y))
conn2.line.color.rgb = rgb(MINT)
conn2.line.width = Pt(2.25)
add_arrow_connector(s38, 10.95, converge_y, 10.95, converge_y + 0.35, color=TEAL)

final_y = converge_y + 0.35
add_rect(s38, 5.5, final_y, 11, 0.95, TEAL, rounded=True, radius=0.1)
add_text(s38, 5.9, final_y, 10.2, 0.95, [("Etapa final consolidada", 15.5, True, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)

add_text(s38, 1.147, final_y + 1.25, 17.7, 0.5,
         [("Use quando um processo se divide em execuções paralelas que depois retornam a um "
           "ponto único — não apenas para fluxos lineares.", 12.5, True, TEAL, True)])

set_footer_label(s38, "TEMPLATE — Fluxograma com bifurcação e convergência")

notes3 = s38.notes_slide
notes3.notes_text_frame.text = ("TEMPLATE - Fluxograma de processo com bifurcação e convergência. "
                                 "Diferente do fluxograma linear (slide 21): aqui uma etapa se ramifica "
                                 "em processos paralelos (setas para baixo a partir de uma linha "
                                 "horizontal) e depois converge de volta a uma etapa única. Use para "
                                 "processos operacionais com sub-fluxos simultâneos, como desmonte, "
                                 "triagem ou qualquer etapa que 'bifurca aqui, converge ali'.")

prs.save(r"C:\Users\andre.eiras\Desktop\Apresentações Ek\build\master_stage2.pptx")
print("saved")
